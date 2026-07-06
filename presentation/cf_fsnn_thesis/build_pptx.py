"""build_pptx.py — faithful PowerPoint export of the reveal.js deck.

Quarto's native `--to pptx` drops ALL figures/animations and the whole reveal design
(twocol/cards/dark theme), because the deck is built on reveal.js/HTML features pptx
cannot express. This builder instead renders each slide to a full-slide image (design
preserved EXACTLY), overlays the real animated GIFs on animated slides (PowerPoint plays
GIFs in slideshow), and carries the speaker notes across.

Deterministic given the rendered deck: non-animated slides are pixel-stable; on animated
slides the static screenshot frame is covered by the overlaid GIF, so the visible result
is stable. Run AFTER `quarto render slides_slim.qmd`.

Usage:  python build_pptx.py   ->   _output/slides_slim.pptx
Requires: playwright (chromium), python-pptx, Pillow.
"""
import pathlib, tempfile, sys
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu

HERE = pathlib.Path(__file__).resolve().parent
DECK = HERE / "_output" / "deck_slim.html"
OUT = HERE / "_output" / "slides_slim.pptx"
GIF_DIR = HERE / "assets" / "manim"

VW, VH = 1280, 720          # reveal viewport
EMU_PER_PX = 9525           # 1280 px -> 13.333in slide (96 dpi); uniform x & y
SLIDE_W, SLIDE_H = VW * EMU_PER_PX, VH * EMU_PER_PX


def capture():
    """Screenshot every slide (2x) + record animated-gif src/bbox + notes. Returns list of
    dicts: {png, gif, bbox, notes}."""
    tmp = pathlib.Path(tempfile.mkdtemp(prefix="pptx_slides_"))
    slides = []
    with sync_playwright() as p:
        b = p.chromium.launch()
        pg = b.new_page(viewport={"width": VW, "height": VH}, device_scale_factor=2)
        pg.goto(DECK.as_uri(), wait_until="networkidle"); pg.wait_for_timeout(1400)
        pg.evaluate("()=>Reveal.slide(0,0)"); pg.wait_for_timeout(300)
        idx = 0
        for _ in range(VH):  # safe upper bound on slide count
            pg.wait_for_timeout(140)
            png = tmp / f"s{idx:02d}.png"
            pg.screenshot(path=str(png))
            info = pg.evaluate("""()=>{
                const s=Reveal.getCurrentSlide();
                const im=s.querySelector('.vis img, img');
                let gif=null;
                if(im && im.src.toLowerCase().endsWith('.gif')){
                    const r=im.getBoundingClientRect();
                    gif={name:im.src.split('/').pop(),
                         x:r.left, y:r.top, w:r.width, h:r.height};
                }
                const note=s.querySelector('aside.notes, .notes');
                return {gif, notes: note ? note.innerText : ''};
            }""")
            slides.append({"png": png, "gif": info["gif"], "notes": info["notes"]})
            if pg.evaluate("()=>Reveal.isLastSlide()"): break
            pg.evaluate("()=>Reveal.next()"); idx += 1
        b.close()
    return slides


def build(slides):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W); prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]
    n_anim = 0
    for s in slides:
        slide = prs.slides.add_slide(blank)
        # full-slide background image (exact design)
        slide.shapes.add_picture(str(s["png"]), 0, 0, width=Emu(SLIDE_W), height=Emu(SLIDE_H))
        # overlay the real GIF on animated slides so it plays in slideshow
        g = s["gif"]
        if g:
            gp = GIF_DIR / g["name"]
            if gp.exists():
                slide.shapes.add_picture(str(gp),
                    Emu(int(g["x"] * EMU_PER_PX)), Emu(int(g["y"] * EMU_PER_PX)),
                    Emu(int(g["w"] * EMU_PER_PX)), Emu(int(g["h"] * EMU_PER_PX)))
                n_anim += 1
        # speaker notes
        if s["notes"].strip():
            slide.notes_slide.notes_text_frame.text = s["notes"].strip()
    prs.save(str(OUT))
    return len(slides), n_anim


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    if not DECK.exists():
        sys.exit(f"deck not found: {DECK} — run `quarto render slides_slim.qmd` first")
    slides = capture()
    n, na = build(slides)
    print(f"OK  {OUT}  —  {n} slides, {na} animated GIFs overlaid, notes carried")
